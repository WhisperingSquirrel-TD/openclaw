#!/usr/bin/env python3
"""
sharepoint_binary_extractor.py — text (and optional image) extraction
from SharePoint binary files for the OpenClaw cache system.

Imported by:
  • sharepoint_cache_poller.py     (automatic, every 15 min)
  • sharepoint_queue_processor.py  (on-demand via read_binary queue entry)

Supported formats
-----------------
  .docx    python-docx      pip install python-docx
  .pdf     pdfminer.six     pip install pdfminer.six
  .pptx    python-pptx      pip install python-pptx
  .msg     extract-msg      pip install extract-msg

Each library is imported lazily — a missing library for one format does not
affect other formats. The returned string describes the failure clearly so
L1 always gets useful output rather than a silent empty file.

Images (nice-to-have)
---------------------
Pass image_dir=Path(...) to save embedded images alongside the extracted
text. Images are saved as image001.png / slide001_01.jpg etc. and referenced
in the returned markdown. Callers that only need text pass image_dir=None.

PDF image extraction requires pymupdf (pip install pymupdf) — if not
installed, text-only extraction still works fine.
"""

from __future__ import annotations

import io
import sys
from pathlib import Path

EXTRACTABLE_EXTENSIONS: frozenset[str] = frozenset({".docx", ".pdf", ".pptx", ".msg"})


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def extract_text(
    filename: str,
    raw_bytes: bytes,
    *,
    image_dir: Path | None = None,
) -> str:
    """Extract text (and optionally images) from a binary SharePoint file.

    Returns a markdown string. Never raises — on failure returns an error
    note that is itself valid markdown, so callers can always write the result
    to the cache regardless of outcome.

    Args:
        filename:   Original filename, used only to detect format via suffix.
        raw_bytes:  Raw file content as bytes.
        image_dir:  If provided, embedded images are saved here and referenced
                    in the returned markdown as relative links. Pass None for
                    text-only extraction.
    """
    ext = Path(filename).suffix.lower()
    try:
        if ext == ".docx":
            return _extract_docx(raw_bytes, image_dir=image_dir)
        if ext == ".pdf":
            return _extract_pdf(raw_bytes, image_dir=image_dir)
        if ext == ".pptx":
            return _extract_pptx(raw_bytes, image_dir=image_dir)
        if ext == ".msg":
            return _extract_msg(raw_bytes)
        return f"_Unsupported file type: `{ext}` — cannot extract text._"
    except ImportError as e:
        pkg = _required_package(ext)
        return (
            f"_Text extraction unavailable: `{pkg}` is not installed._\n\n"
            f"Install with: `pip3 install --break-system-packages {pkg}`\n\n"
            f"_(ImportError: {e})_"
        )
    except Exception as e:
        return f"_Text extraction failed for `{filename}`: {e}_"


def _required_package(ext: str) -> str:
    return {
        ".docx": "python-docx",
        ".pdf":  "pdfminer.six",
        ".pptx": "python-pptx",
        ".msg":  "extract-msg",
    }.get(ext, "unknown")


# ---------------------------------------------------------------------------
# .docx — python-docx
# ---------------------------------------------------------------------------

def _extract_docx(raw_bytes: bytes, *, image_dir: Path | None = None) -> str:
    from docx import Document  # python-docx

    doc   = Document(io.BytesIO(raw_bytes))
    parts: list[str] = []
    img_counter = [0]

    def _save_image(blob: bytes, content_type: str) -> str | None:
        if image_dir is None:
            return None
        image_dir.mkdir(parents=True, exist_ok=True)
        img_counter[0] += 1
        ext_map = {
            "image/png":  ".png",
            "image/jpeg": ".jpg",
            "image/gif":  ".gif",
            "image/svg+xml": ".svg",
        }
        ext     = ext_map.get(content_type, ".bin")
        fname   = f"image{img_counter[0]:03d}{ext}"
        (image_dir / fname).write_bytes(blob)
        return fname

    # Walk body XML in document order so heading / paragraph / table
    # sequence is faithfully preserved.
    for block in doc.element.body:
        local = block.tag.split("}")[-1] if "}" in block.tag else block.tag

        if local == "p":
            from docx.oxml.ns import qn
            from docx.text.paragraph import Paragraph

            para = Paragraph(block, doc)

            # Extract embedded images in this paragraph
            img_refs: list[str] = []
            for blip in block.iter(qn("a:blip")):
                r_embed = blip.get(qn("r:embed"))
                if r_embed:
                    try:
                        img_part = doc.part.related_parts.get(r_embed)
                        if img_part:
                            fname = _save_image(img_part.blob, img_part.content_type)
                            if fname:
                                img_refs.append(f"![image](./{fname})")
                    except Exception:
                        pass

            text  = para.text.strip()
            style = para.style.name if para.style else ""

            if style.startswith("Heading"):
                level  = "".join(c for c in style if c.isdigit()) or "1"
                prefix = "#" * min(int(level), 6)
                if text:
                    parts.append(f"{prefix} {text}")
            else:
                if text:
                    parts.append(text)

            parts.extend(img_refs)

        elif local == "tbl":
            from docx.table import Table

            tbl  = Table(block, doc)
            rows: list[str] = []
            for i, row in enumerate(tbl.rows):
                cells = [c.text.strip().replace("\n", " ") for c in row.cells]
                rows.append("| " + " | ".join(cells) + " |")
                if i == 0:
                    rows.append("|" + "|".join(["---"] * len(cells)) + "|")
            parts.extend(rows)

    return "\n\n".join(p for p in parts if p) or "_No text content found in document._"


# ---------------------------------------------------------------------------
# .pdf — pdfminer.six (text); pymupdf (images, optional)
# ---------------------------------------------------------------------------

def _extract_pdf(raw_bytes: bytes, *, image_dir: Path | None = None) -> str:
    from pdfminer.high_level import extract_text_to_fp  # pdfminer.six
    from pdfminer.layout import LAParams

    out = io.StringIO()
    extract_text_to_fp(
        io.BytesIO(raw_bytes), out,
        laparams=LAParams(),
        output_type="text",
        codec="utf-8",
    )
    text = out.getvalue().strip()

    if image_dir is not None:
        try:
            import fitz  # pymupdf — optional; text works fine without it
            pdf_doc  = fitz.open(stream=raw_bytes, filetype="pdf")
            img_refs: list[str] = []
            for page_num, page in enumerate(pdf_doc, 1):
                for img_idx, img_info in enumerate(page.get_images(), 1):
                    xref     = img_info[0]
                    base_img = pdf_doc.extract_image(xref)
                    if not base_img:
                        continue
                    ext   = base_img.get("ext", "png")
                    fname = f"page{page_num:03d}_img{img_idx:02d}.{ext}"
                    image_dir.mkdir(parents=True, exist_ok=True)
                    (image_dir / fname).write_bytes(base_img["image"])
                    img_refs.append(f"![image p{page_num}](./{fname})")
            if img_refs:
                text += "\n\n---\n\n**Extracted images:**\n\n" + "\n".join(img_refs)
        except ImportError:
            pass  # pymupdf not installed — text-only is fine

    return text or "_No text content extracted from PDF._"


# ---------------------------------------------------------------------------
# .pptx — python-pptx
# ---------------------------------------------------------------------------

def _extract_pptx(raw_bytes: bytes, *, image_dir: Path | None = None) -> str:
    from pptx import Presentation  # python-pptx

    prs   = Presentation(io.BytesIO(raw_bytes))
    parts: list[str] = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_parts: list[str] = [f"## Slide {slide_num}"]

        title_text = ""
        if slide.shapes.title and slide.shapes.title.text.strip():
            title_text = slide.shapes.title.text.strip()
            slide_parts.append(f"### {title_text}")

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text and text != title_text:
                        slide_parts.append(text)

            if shape.has_table:
                rows: list[str] = []
                for i, row in enumerate(shape.table.rows):
                    cells = [c.text.strip().replace("\n", " ") for c in row.cells]
                    rows.append("| " + " | ".join(cells) + " |")
                    if i == 0:
                        rows.append("|" + "|".join(["---"] * len(cells)) + "|")
                slide_parts.extend(rows)

            # MSO_SHAPE_TYPE.PICTURE == 13
            if image_dir is not None and shape.shape_type == 13:
                try:
                    img   = shape.image
                    ext_m = {"image/png": ".png", "image/jpeg": ".jpg", "image/gif": ".gif"}
                    ext   = ext_m.get(img.content_type, ".bin")
                    fname = f"slide{slide_num:03d}_{shape.shape_id}{ext}"
                    image_dir.mkdir(parents=True, exist_ok=True)
                    (image_dir / fname).write_bytes(img.blob)
                    slide_parts.append(f"![slide image](./{fname})")
                except Exception:
                    pass

        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_parts.append(f"**Notes:** {notes}")

        parts.append("\n\n".join(p for p in slide_parts if p))

    return "\n\n---\n\n".join(parts) or "_No text content found in presentation._"


# ---------------------------------------------------------------------------
# .msg — extract-msg
# ---------------------------------------------------------------------------

def _extract_msg(raw_bytes: bytes) -> str:
    import extract_msg as em  # extract-msg

    with em.openMsg(io.BytesIO(raw_bytes)) as msg:
        lines = [
            f"**From:** {msg.sender or '(unknown)'}",
            f"**To:** {msg.to or '(unknown)'}",
            f"**CC:** {msg.cc or ''}",
            f"**Date:** {msg.date or ''}",
            f"**Subject:** {msg.subject or '(no subject)'}",
            "",
            "---",
            "",
        ]

        body = (msg.body or "").strip()
        if body:
            lines.append(body)
        else:
            lines.append("_(No plain-text body)_")

        if msg.attachments:
            lines.append("")
            lines.append(f"**Attachments ({len(msg.attachments)}):**")
            for att in msg.attachments:
                name = (
                    getattr(att, "longFilename", None)
                    or getattr(att, "shortFilename", None)
                    or "(unnamed)"
                )
                lines.append(f"- {name}")

    return "\n".join(lines)


if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(description="Extract text from a binary SharePoint file.")
    parser.add_argument("file", help="Path to the file to extract")
    parser.add_argument("--images", help="Directory to save extracted images (optional)")
    args = parser.parse_args()

    path      = Path(args.file)
    image_dir = Path(args.images) if args.images else None
    result    = extract_text(path.name, path.read_bytes(), image_dir=image_dir)
    print(result)
